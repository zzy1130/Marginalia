"""Screen capture and active document detection for macOS."""

import os
import re
import subprocess
import tempfile
from typing import Optional


# Document viewer apps we look for
_DOC_APPS = {"Preview", "预览", "Microsoft PowerPoint", "Keynote", "Adobe Acrobat",
             "Adobe Acrobat Reader", "Skim", "PDF Expert", "WPS Office"}

# File extensions we care about
_DOC_EXTENSIONS = (".pdf", ".pptx", ".ppt", ".key")


def get_document_window_info() -> dict:
    """Find the most relevant document viewer window, skipping our own app.

    Scans all on-screen windows and finds one belonging to a known
    document viewer app (Preview, PowerPoint, etc.).
    Falls back to the most recent non-self window with a document title.
    """
    from Quartz import (
        CGWindowListCopyWindowInfo,
        kCGWindowListExcludeDesktopElements,
        kCGWindowListOptionOnScreenOnly,
        kCGNullWindowID,
    )

    my_pid = os.getpid()

    window_list = CGWindowListCopyWindowInfo(
        kCGWindowListOptionOnScreenOnly | kCGWindowListExcludeDesktopElements,
        kCGNullWindowID,
    )

    best = None
    fallback = None

    for window in window_list:
        pid = window.get("kCGWindowOwnerPID", 0)
        if pid == my_pid:
            continue
        if window.get("kCGWindowLayer", -1) != 0:
            continue

        app_name = window.get("kCGWindowOwnerName", "")
        title = window.get("kCGWindowName", "") or ""
        window_id = window.get("kCGWindowNumber")

        # Priority 1: known document viewer app
        if app_name in _DOC_APPS:
            if best is None:
                best = {
                    "app_name": app_name,
                    "title": title,
                    "window_id": window_id,
                    "pid": pid,
                }

        # Priority 2: any window whose title contains a doc filename
        if fallback is None and any(ext in title.lower() for ext in _DOC_EXTENSIONS):
            fallback = {
                "app_name": app_name,
                "title": title,
                "window_id": window_id,
                "pid": pid,
            }

    if best:
        return best
    if fallback:
        return fallback

    return {
        "app_name": "",
        "title": "",
        "window_id": None,
        "pid": 0,
    }


def capture_window(window_id: int, output_path: Optional[str] = None) -> str:
    """Capture a specific window by its macOS window ID."""
    if output_path is None:
        output_path = os.path.join(tempfile.gettempdir(), "marginalia_capture.png")
    subprocess.run(
        ["screencapture", "-l", str(window_id), "-o", output_path],
        check=True,
    )
    return output_path


def parse_page_info(app_name: str, title: str) -> dict:
    """Parse filename and page number from window title."""
    info = {"filename": None, "current_page": None, "total_pages": None}

    # Extract filename (PDF/PPT/KEY)
    match = re.match(r"^(.+?\.(?:pdf|pptx?|key))", title, re.IGNORECASE)
    if match:
        info["filename"] = match.group(1).strip()

    if app_name in ("Preview", "\u9884\u89c8"):
        # English: "file.pdf (page 3 of 10)"
        m = re.search(r"page\s+(\d+)\s+of\s+(\d+)", title, re.IGNORECASE)
        if m:
            info["current_page"] = int(m.group(1))
            info["total_pages"] = int(m.group(2))
        else:
            # Chinese: "第 3 页" / "共 10 页"
            m = re.search(r"\u7b2c\s*(\d+)\s*\u9875", title)
            if m:
                info["current_page"] = int(m.group(1))
            m = re.search(r"\u5171\s*(\d+)\s*\u9875", title)
            if m:
                info["total_pages"] = int(m.group(1))

    elif "PowerPoint" in app_name:
        try:
            result = subprocess.run(
                [
                    "osascript", "-e",
                    'tell application "Microsoft PowerPoint" to get '
                    "slide index of slide range of selection of active window",
                ],
                capture_output=True, text=True, timeout=5,
            )
            if result.returncode == 0 and result.stdout.strip().isdigit():
                info["current_page"] = int(result.stdout.strip())
        except Exception:
            pass

    elif "Keynote" in app_name:
        try:
            result = subprocess.run(
                [
                    "osascript", "-e",
                    'tell application "Keynote" to get slide number '
                    "of current slide of front document",
                ],
                capture_output=True, text=True, timeout=5,
            )
            if result.returncode == 0 and result.stdout.strip().isdigit():
                info["current_page"] = int(result.stdout.strip())
        except Exception:
            pass

    return info


def find_file_path(app_name: str, pid: int, filename: Optional[str] = None) -> Optional[str]:
    """Try to find the actual file path of the open document via lsof, then Spotlight."""
    # Method 1: lsof
    if pid > 0:
        try:
            result = subprocess.run(
                ["lsof", "-p", str(pid)],
                capture_output=True, text=True, timeout=5,
            )
            for line in result.stdout.split("\n"):
                lower = line.lower()
                if not any(ext in lower for ext in _DOC_EXTENSIONS):
                    continue
                idx = line.find("/")
                if idx == -1:
                    continue
                path = line[idx:].strip()
                if os.path.exists(path):
                    if filename and filename.lower() in os.path.basename(path).lower():
                        return path
                    elif not filename:
                        return path
        except Exception:
            pass

    # Method 2: Spotlight (mdfind) — handles sandboxed apps like Preview
    if filename:
        try:
            result = subprocess.run(
                ["mdfind", "-name", filename],
                capture_output=True, text=True, timeout=5,
            )
            for line in result.stdout.strip().split("\n"):
                path = line.strip()
                if path and os.path.exists(path) and any(path.lower().endswith(ext) for ext in _DOC_EXTENSIONS):
                    return path
        except Exception:
            pass

    return None


def extract_page_text(file_path: str, page_number: int) -> Optional[str]:
    """Extract text content from a specific page of a PDF or PPTX file."""
    if not file_path or not os.path.exists(file_path):
        return None

    lower = file_path.lower()

    if lower.endswith(".pdf"):
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(file_path)
            if 1 <= page_number <= len(doc):
                page = doc[page_number - 1]
                text = page.get_text()
                doc.close()
                return text.strip() if text.strip() else None
            doc.close()
        except Exception as e:
            print(f"[DEBUG] PDF text extraction failed: {e}", flush=True)

    elif lower.endswith((".pptx", ".ppt")):
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            if 1 <= page_number <= len(prs.slides):
                slide = prs.slides[page_number - 1]
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                texts.append(t)
                return "\n".join(texts) if texts else None
        except Exception as e:
            print(f"[DEBUG] PPTX text extraction failed: {e}", flush=True)

    return None


def get_document_context() -> dict:
    """Get full context about the currently viewed document."""
    win = get_document_window_info()
    page = parse_page_info(win["app_name"], win["title"])
    file_path = find_file_path(win["app_name"], win["pid"], page["filename"])

    # Extract text from the current page (default to page 1 if unknown)
    page_text = None
    if file_path:
        page_num = page["current_page"] or 1
        page_text = extract_page_text(file_path, page_num)

    return {
        "app_name": win["app_name"],
        "title": win["title"],
        "window_id": win["window_id"],
        "filename": page["filename"],
        "current_page": page["current_page"],
        "total_pages": page["total_pages"],
        "file_path": file_path,
        "page_text": page_text,
    }
