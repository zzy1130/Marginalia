"""OpenSandbox integration for secure code execution.

Provides SandboxManager that isolates run_code and insert_notes
inside Docker containers via OpenSandbox, with graceful fallback
to cloud-hosted sandbox or restricted local execution.
"""

import asyncio
import logging
import os
import shutil
import subprocess
import tempfile
from datetime import timedelta
from typing import Optional

from dotenv import dotenv_values

logger = logging.getLogger("marginalia.sandbox")

_config = dotenv_values(os.path.join(os.path.dirname(__file__), "..", ".env"))

SANDBOX_IMAGE = "opensandbox/code-interpreter:v1.0.1"
SANDBOX_ENTRYPOINT = ["/opt/opensandbox/code-interpreter.sh"]
LOCAL_SERVER_PORT = 8080
_VENV_PYTHON = os.path.join(os.path.dirname(__file__), "..", ".venv", "bin", "python")

_ALLOWED_EXTENSIONS = (".pdf", ".pptx", ".ppt")


def validate_file_path(file_path: str) -> str:
    """Validate and resolve a file path from LLM tool input.

    Prevents path traversal, symlink attacks, and writes outside home directory.
    Returns the resolved real path.
    """
    real = os.path.realpath(file_path)
    if not os.path.exists(real):
        raise ValueError(f"File not found: {file_path}")
    if not real.lower().endswith(_ALLOWED_EXTENSIONS):
        raise ValueError(f"Unsupported file type: {file_path}")
    home = os.path.expanduser("~")
    if not real.startswith(home + os.sep) and real != home:
        raise ValueError(f"File outside home directory: {file_path}")
    return real


class SandboxManager:
    """Manages OpenSandbox lifecycle for secure tool execution.

    Supports three modes:
    - local: Run opensandbox-server locally with Docker
    - cloud: Use api.opensandbox.io
    - disabled: Fall back to restricted local subprocess
    """

    def __init__(self):
        self._mode = _config.get("SANDBOX_MODE", "auto")
        self._cloud_api_key = _config.get("OPENSANDBOX_API_KEY", "")
        self._cloud_domain = _config.get("OPENSANDBOX_DOMAIN", "api.opensandbox.io")

        self._sandbox = None
        self._interpreter = None
        self._connection_config = None
        self._local_server_proc = None
        self._use_cloud = False
        self._initialized = False
        self._available = False

    async def initialize(self) -> None:
        """Detect Docker, start sandbox server, create warm sandbox."""
        if self._initialized:
            return
        self._initialized = True

        if self._mode == "disabled":
            logger.info("Sandbox disabled by configuration")
            return

        if self._mode in ("auto", "local"):
            if self._is_docker_running():
                try:
                    await self._start_local_server()
                    await self._create_warm_sandbox()
                    self._available = True
                    logger.info("Sandbox initialized with local Docker")
                    return
                except Exception as e:
                    logger.warning("Failed to start local sandbox: %s", e)
                    if self._mode == "local":
                        return

        if self._mode in ("auto", "cloud"):
            if self._cloud_api_key:
                try:
                    self._use_cloud = True
                    await self._create_warm_sandbox()
                    self._available = True
                    logger.info("Sandbox initialized with cloud provider (%s)", self._cloud_domain)
                    return
                except Exception as e:
                    logger.warning("Failed to connect to cloud sandbox: %s", e)
                    self._use_cloud = False

        logger.warning(
            "No sandbox available (mode=%s). Tool execution will use restricted local fallback.",
            self._mode,
        )

    @property
    def is_available(self) -> bool:
        return self._available

    def _is_docker_running(self) -> bool:
        docker_bin = shutil.which("docker")
        if not docker_bin:
            logger.info("Docker not found in PATH")
            return False
        try:
            result = subprocess.run(
                [docker_bin, "info"],
                capture_output=True, timeout=10,
            )
            return result.returncode == 0
        except Exception:
            return False

    async def _start_local_server(self) -> None:
        """Start opensandbox-server locally if not already running."""
        import httpx

        server_url = f"http://127.0.0.1:{LOCAL_SERVER_PORT}"
        try:
            async with httpx.AsyncClient(timeout=3.0) as client:
                r = await client.get(f"{server_url}/health")
                if r.status_code == 200:
                    logger.info("Local opensandbox-server already running")
                    await self._setup_local_connection()
                    return
        except Exception:
            pass

        config_path = os.path.expanduser("~/.sandbox.toml")
        if not os.path.exists(config_path):
            await self._create_default_config(config_path)

        venv_bin = os.path.join(os.path.dirname(__file__), "..", ".venv", "bin", "opensandbox-server")
        opensandbox_bin = venv_bin if os.path.exists(venv_bin) else shutil.which("opensandbox-server")
        if not opensandbox_bin:
            raise RuntimeError(
                "opensandbox-server not found. Install with: uv pip install opensandbox-server"
            )

        self._local_server_proc = subprocess.Popen(
            [opensandbox_bin, "--config", config_path],
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
        )

        for _ in range(30):
            await asyncio.sleep(1)
            try:
                async with httpx.AsyncClient(timeout=3.0) as client:
                    r = await client.get(f"{server_url}/health")
                    if r.status_code == 200:
                        await self._setup_local_connection()
                        return
            except Exception:
                continue

        raise RuntimeError("opensandbox-server failed to start within 30s")

    async def _create_default_config(self, path: str) -> None:
        config = f"""\
[server]
host = "127.0.0.1"
port = {LOCAL_SERVER_PORT}
log_level = "WARNING"

[runtime]
type = "docker"
execd_image = "opensandbox/execd:v1.0.6"

[egress]
image = "opensandbox/egress:v1.0.1"

[docker]
network_mode = "bridge"
drop_capabilities = ["AUDIT_WRITE", "MKNOD", "NET_ADMIN", "NET_RAW", "SYS_ADMIN", "SYS_MODULE", "SYS_PTRACE", "SYS_TIME", "SYS_TTY_CONFIG"]
no_new_privileges = true
pids_limit = 256
"""
        with open(path, "w") as f:
            f.write(config)

    async def _setup_local_connection(self) -> None:
        from opensandbox.config import ConnectionConfig

        self._connection_config = ConnectionConfig(
            domain=f"127.0.0.1:{LOCAL_SERVER_PORT}",
            api_key="",
            protocol="http",
            request_timeout=timedelta(seconds=60),
        )

    async def _create_warm_sandbox(self) -> None:
        """Pre-create a sandbox container so tool calls don't pay cold-start cost."""
        from opensandbox import Sandbox
        from opensandbox.models.sandboxes import NetworkPolicy, NetworkRule
        from code_interpreter import CodeInterpreter

        if self._use_cloud:
            from opensandbox.config import ConnectionConfig

            self._connection_config = ConnectionConfig(
                domain=self._cloud_domain,
                api_key=self._cloud_api_key,
                request_timeout=timedelta(seconds=60),
            )

        self._sandbox = await Sandbox.create(
            SANDBOX_IMAGE,
            connection_config=self._connection_config,
            entrypoint=SANDBOX_ENTRYPOINT,
            env={"PYTHON_VERSION": "3.11"},
            timeout=timedelta(minutes=30),
            resource={"cpu": "1", "memory": "1Gi"},
            network_policy=NetworkPolicy(
                defaultAction="deny",
                egress=[
                    NetworkRule(action="allow", target="pypi.org"),
                    NetworkRule(action="allow", target="*.python.org"),
                ],
            ),
        )

        self._interpreter = await CodeInterpreter.create(sandbox=self._sandbox)
        logger.info("Warm sandbox created: %s", self._sandbox.id)

    async def _ensure_sandbox(self) -> None:
        """Ensure sandbox is alive, recreate if expired."""
        if self._sandbox is None:
            await self._create_warm_sandbox()
            return

        try:
            info = await self._sandbox.get_info()
            if info.status.state not in ("Running", "RUNNING"):
                logger.info("Sandbox expired (state=%s), recreating", info.status.state)
                await self._create_warm_sandbox()
        except Exception:
            logger.info("Sandbox unreachable, recreating")
            await self._create_warm_sandbox()

    async def run_code(self, code: str) -> str:
        """Execute Python code inside the sandbox.

        Returns combined stdout/stderr/result output.
        """
        if not self._available:
            return self._run_code_fallback(code)

        try:
            await self._ensure_sandbox()
            from code_interpreter import SupportedLanguage

            result = await self._interpreter.codes.run(
                code, language=SupportedLanguage.PYTHON,
            )

            output_parts = []
            if result.logs and result.logs.stdout:
                output_parts.append("".join(m.text for m in result.logs.stdout))
            if result.logs and result.logs.stderr:
                stderr = "".join(m.text for m in result.logs.stderr)
                if stderr.strip():
                    output_parts.append(f"STDERR:\n{stderr}")
            if result.result:
                output_parts.append("".join(r.text for r in result.result))

            return "\n".join(output_parts) or "(no output)"

        except Exception as e:
            logger.error("Sandbox run_code failed: %s, falling back to local", e)
            return self._run_code_fallback(code)

    def _run_code_fallback(self, code: str) -> str:
        """Restricted local subprocess fallback when no sandbox is available.

        Blocks dangerous modules by prepending import guards.
        """
        blocked_modules = [
            "os", "subprocess", "shutil", "socket", "http", "urllib",
            "ftplib", "smtplib", "ctypes", "importlib", "sys",
            "pathlib", "glob", "tempfile", "signal", "multiprocessing",
        ]
        guard = "import builtins as _b\n"
        guard += "_orig_import = _b.__import__\n"
        guard += f"_BLOCKED = {set(blocked_modules)}\n"
        guard += (
            "def _safe_import(name, *a, **kw):\n"
            "    if name.split('.')[0] in _BLOCKED:\n"
            "        raise ImportError(f'Module {name} is blocked in restricted mode')\n"
            "    return _orig_import(name, *a, **kw)\n"
            "_b.__import__ = _safe_import\n"
        )
        restricted_code = guard + code

        env = os.environ.copy()
        env["DYLD_LIBRARY_PATH"] = "/opt/homebrew/lib:" + env.get("DYLD_LIBRARY_PATH", "")
        env["DYLD_FALLBACK_LIBRARY_PATH"] = "/opt/homebrew/lib"

        try:
            result = subprocess.run(
                [os.path.abspath(_VENV_PYTHON), "-c", restricted_code],
                capture_output=True,
                text=True,
                timeout=30,
                cwd=tempfile.gettempdir(),
                env=env,
            )
            output = ""
            if result.stdout:
                output += result.stdout
            if result.stderr:
                output += f"\nSTDERR:\n{result.stderr}"
            if result.returncode != 0:
                output += f"\nExit code: {result.returncode}"
            return output or "(no output)"
        except subprocess.TimeoutExpired:
            return "Error: Code execution timed out (30s limit)."
        except Exception as e:
            return f"Error running code: {e}"

    async def run_insert_notes(
        self, markdown_content: str, file_path: str, page_index: int,
    ) -> str:
        """Execute insert_notes locally with path validation.

        insert_notes runs our own controlled code (not arbitrary LLM code),
        so it's safe to execute locally. The sandbox is only needed for
        run_code which executes arbitrary LLM-generated code.

        Security is provided by validate_file_path() which prevents path
        traversal, symlink attacks, and writes outside the home directory.
        """
        validated_path = validate_file_path(file_path)
        return self._insert_notes_fallback(markdown_content, validated_path, page_index)

    def _insert_notes_fallback(
        self, markdown_content: str, file_path: str, page_index: int,
    ) -> str:
        """Local fallback for insert_notes with the f-string injection fixed."""
        import fitz
        import markdown as md
        import re

        formula_counter = [0]
        formula_dir = os.path.join(tempfile.gettempdir(), "marginalia_formulas")
        os.makedirs(formula_dir, exist_ok=True)

        def render_latex(match):
            formula = match.group(1)
            is_display = match.group(0).startswith("$$")
            formula_counter[0] += 1
            img_path = os.path.join(formula_dir, f"f{formula_counter[0]}.png")
            try:
                import matplotlib
                matplotlib.use("Agg")
                import matplotlib.pyplot as plt
                from PIL import Image
                fig, ax = plt.subplots(figsize=(0.01, 0.01))
                ax.axis("off")
                fontsize = 14 if is_display else 11
                ax.text(0, 0, f"${formula}$", fontsize=fontsize, color="#222",
                        math_fontfamily="cm")
                fig.savefig(img_path, bbox_inches="tight", pad_inches=0.02,
                            dpi=300, transparent=True)
                plt.close(fig)
                with Image.open(img_path) as im:
                    pw, ph = im.size
                target_h = 18 if is_display else 14
                scale = target_h / max(ph, 1)
                w = int(pw * scale)
                h = target_h
            except Exception:
                return match.group(0)
            img_name = os.path.basename(img_path)
            if is_display:
                return f'<p style="text-align:center;margin:8px 0;"><img src="{img_name}" width="{w}" height="{h}"></p>'
            else:
                return f'<img src="{img_name}" width="{w}" height="{h}" style="vertical-align:text-bottom;">'

        mc = re.sub(r'\$\$(.+?)\$\$', render_latex, markdown_content, flags=re.DOTALL)
        mc = re.sub(r'\$(.+?)\$', render_latex, mc)

        body = md.markdown(mc, extensions=["fenced_code", "tables", "nl2br"])
        body += (
            '<p style="margin-top: 24px; padding-top: 8px; border-top: 1px solid #ddd; '
            'color: #aaa; font-size: 8px; text-align: right;">Marginalia Notes</p>'
        )

        from core.agent import _NOTES_CSS

        MEDIABOX = fitz.paper_rect("a4")
        WHERE = MEDIABOX + (50, 50, -50, -40)

        pdf_path = os.path.join(tempfile.gettempdir(), "marginalia_notes.pdf")
        writer = fitz.DocumentWriter(pdf_path)
        story = fitz.Story(html=body, user_css=_NOTES_CSS, archive=formula_dir)

        more = True
        while more:
            dev = writer.begin_page(MEDIABOX)
            more, _ = story.place(WHERE)
            story.draw(dev)
            writer.end_page()
        writer.close()

        lower = file_path.lower()
        if lower.endswith(".pdf"):
            doc = fitz.open(file_path)
            tmp_doc = fitz.open(pdf_path)
            doc.insert_pdf(tmp_doc, from_page=0, to_page=tmp_doc.page_count - 1,
                           start_at=page_index + 1)
            doc.saveIncr()
            doc.close()
            tmp_doc.close()
            return "Notes inserted into PDF successfully!"
        elif lower.endswith((".pptx", ".ppt")):
            from pptx import Presentation
            from pptx.util import Inches, Pt
            prs = Presentation(file_path)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = "Study Notes"
            p.font.size = Pt(24)
            p.font.bold = True
            prs.save(file_path)
            return "Notes slide added to PPTX successfully!"
        else:
            return f"Unsupported file type: {file_path}"

    async def shutdown(self) -> None:
        """Clean up sandbox and local server."""
        if self._sandbox:
            try:
                await self._sandbox.kill()
            except Exception:
                pass
            self._sandbox = None
            self._interpreter = None

        if self._local_server_proc:
            self._local_server_proc.terminate()
            try:
                self._local_server_proc.wait(timeout=5)
            except subprocess.TimeoutExpired:
                self._local_server_proc.kill()
            self._local_server_proc = None

        logger.info("Sandbox manager shut down")


_manager: Optional[SandboxManager] = None


async def init_sandbox_manager() -> SandboxManager:
    """Initialize the global sandbox manager singleton."""
    global _manager
    if _manager is None:
        _manager = SandboxManager()
        await _manager.initialize()
    return _manager


def get_sandbox_manager() -> SandboxManager:
    """Get the global sandbox manager. Must be initialized first."""
    if _manager is None:
        raise RuntimeError("SandboxManager not initialized. Call init_sandbox_manager() first.")
    return _manager
