"""Insert formatted study notes into PDF/PPTX files.

Helper script for the insert-notes skill.
Called by the LLM via run_code.
"""

import fitz
import markdown as md
import re
import os
import sys
import tempfile


NOTES_CSS = """\
body {
  font-family: "PingFang SC", "Hiragino Sans GB", "Noto Sans CJK SC", sans-serif;
  font-size: 10px; color: #222; line-height: 1.6;
}
h1 {
  font-size: 17px; font-weight: bold; color: #2E4226;
  border-bottom: 2px solid #2E4226; padding-bottom: 6px;
  margin-bottom: 10px;
}
h2 {
  font-size: 12px; font-weight: bold; color: #2E4226;
  margin-top: 12px; margin-bottom: 5px; padding-bottom: 3px;
  border-bottom: 1px solid #d4ddd0;
}
h3 { font-size: 11px; font-weight: bold; color: #55644A; margin-top: 8px; margin-bottom: 4px; }
p { margin-top: 3px; margin-bottom: 3px; }
ul, ol { padding-left: 18px; margin-top: 4px; margin-bottom: 4px; }
li { margin-top: 1px; margin-bottom: 1px; }
b, strong { font-weight: bold; color: #1a3a10; }
code {
  background-color: #f0f0f0; padding: 1px 4px;
  font-family: monospace; font-size: 9px;
}
pre {
  background-color: #f7f8f6; padding: 10px 12px;
  border: 1px solid #e4e8e2; font-size: 9px;
  font-family: monospace; line-height: 1.5;
}
blockquote {
  border-left: 3px solid #2E4226; padding: 5px 10px;
  background-color: #f2f6ef; color: #333; font-style: italic;
  margin-top: 6px; margin-bottom: 6px;
}
table { border-collapse: collapse; margin-top: 6px; margin-bottom: 6px; }
th {
  font-weight: bold; text-align: left; padding: 4px 8px;
  border: 1px solid #d4ddd0; font-size: 9px; color: #2E4226;
  background-color: #eef3eb;
}
td { padding: 3px 8px; border: 1px solid #dde3da; font-size: 9px; }
hr { border: none; border-top: 1px solid #e0e0e0; margin-top: 8px; margin-bottom: 8px; }
"""


def _render_latex(match, formula_counter, formula_dir):
    formula = match.group(1)
    is_display = match.group(0).startswith("$$")
    formula_counter[0] += 1
    img_path = os.path.join(formula_dir, f"f{formula_counter[0]}.png")
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        from PIL import Image
        fig, ax = plt.subplots(figsize=(0.01, 0.01))
        ax.axis("off")
        fontsize = 14 if is_display else 11
        ax.text(0, 0, f"${formula}$", fontsize=fontsize, color="#222",
                math_fontfamily="cm")
        fig.savefig(img_path, bbox_inches="tight", pad_inches=0.02,
                    dpi=300, transparent=True)
        plt.close(fig)
        with Image.open(img_path) as im:
            pw, ph = im.size
        target_h = 18 if is_display else 14
        scale = target_h / max(ph, 1)
        w = int(pw * scale)
        h = target_h
    except Exception:
        return match.group(0)
    img_name = os.path.basename(img_path)
    if is_display:
        return f'<p style="text-align:center;margin:8px 0;"><img src="{img_name}" width="{w}" height="{h}"></p>'
    else:
        return f'<img src="{img_name}" width="{w}" height="{h}" style="vertical-align:text-bottom;">'


def insert_notes(markdown_content: str, file_path: str, page_number: int) -> str:
    """page_number is 1-based (same as what the user sees in their PDF viewer)."""
    page_index = max(0, page_number - 1)
    print(f"[insert_notes] file={file_path}, page_number={page_number}, page_index={page_index}, content_len={len(markdown_content)}", file=sys.stderr)

    if not os.path.exists(file_path):
        return f"Error: file not found: {file_path}"

    formula_counter = [0]
    formula_dir = os.path.join(tempfile.gettempdir(), "marginalia_formulas")
    os.makedirs(formula_dir, exist_ok=True)

    def latex_replacer(match):
        return _render_latex(match, formula_counter, formula_dir)

    mc = re.sub(r'\$\$(.+?)\$\$', latex_replacer, markdown_content, flags=re.DOTALL)
    mc = re.sub(r'\$(.+?)\$', latex_replacer, mc)

    body = md.markdown(mc, extensions=["fenced_code", "tables", "nl2br"])
    body += (
        '<p style="margin-top: 24px; padding-top: 8px; border-top: 1px solid #ddd; '
        'color: #aaa; font-size: 8px; text-align: right;">Marginalia Notes</p>'
    )

    MEDIABOX = fitz.paper_rect("a4")
    WHERE = MEDIABOX + (50, 50, -50, -40)

    pdf_path = os.path.join(tempfile.gettempdir(), "marginalia_notes.pdf")
    writer = fitz.DocumentWriter(pdf_path)

    archive = fitz.Archive()
    for d in [formula_dir, "/System/Library/Fonts", "/System/Library/Fonts/Supplemental", "/Library/Fonts"]:
        if os.path.isdir(d):
            archive.add(d)

    story = fitz.Story(html=body, user_css=NOTES_CSS, archive=archive)

    more = True
    while more:
        dev = writer.begin_page(MEDIABOX)
        more, _ = story.place(WHERE)
        story.draw(dev)
        writer.end_page()
    writer.close()

    lower = file_path.lower()
    if lower.endswith(".pdf"):
        doc = fitz.open(file_path)
        tmp_doc = fitz.open(pdf_path)
        doc.insert_pdf(tmp_doc, from_page=0, to_page=0,
                       start_at=page_index + 1)
        doc.saveIncr()
        doc.close()
        tmp_doc.close()
        return "Notes inserted into PDF successfully!"
    elif lower.endswith((".pptx", ".ppt")):
        from pptx import Presentation
        from pptx.util import Inches, Pt
        prs = Presentation(file_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "Study Notes"
        p.font.size = Pt(24)
        p.font.bold = True
        prs.save(file_path)
        return "Notes slide added to PPTX successfully!"
    else:
        return f"Unsupported file type: {file_path}"
